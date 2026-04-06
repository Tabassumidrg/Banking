from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define color scheme (Modern Banking - Dark/Blue)
    def set_slide_background(slide, color=RGBColor(10, 25, 47)): # Deep Navy
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_bullet_points(slide, points):
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    # Slide 0: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Nidhi Bank"
    subtitle.text = "A Premium Full-Stack Digital Banking Solution\nNext.js | FastAPI | PostgreSQL"
    
    # Slide 1: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Overview"
    points = [
        "Modern digital banking platform focused on security and user experience.",
        "End-to-end implementation from frontend UI to backend logic.",
        "Built to demonstrate atomic financial transactions and secure auth.",
        "Features a high-end Glassmorphism design system."
    ]
    add_bullet_points(slide, points)

    # Slide 2: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Technology Stack"
    points = [
        "Frontend: Next.js 14, Vanilla CSS, Responsive Design",
        "Backend: FastAPI (Python), RESTful API Architecture",
        "Database: PostgreSQL (Neon Cloud)",
        "Authentication: Bcrypt Hashing, Secure Sign-in Flow",
        "Deployment: Vercel (Frontend), Render (Backend)"
    ]
    add_bullet_points(slide, points)

    # Slide 3: Key Features (Part 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Core Banking Features"
    points = [
        "User Onboarding: Secure signup and profile creation.",
        "Financial Dashboard: Real-time balance and transaction history.",
        "Fund Transfers: Peer-to-peer transfers using receiver email.",
        "User Management: Administrative views for user listing and auditing."
    ]
    add_bullet_points(slide, points)

    # Slide 4: Secure Authentication
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Security Implementation"
    points = [
        "Password Security: No plain-text storage; direct Bcrypt hashing.",
        "CORS Protection: Explicit middleware handling for secure API access.",
        "Data Integrity: Validation at both frontend and backend layers.",
        "Protected Routes: Dashboard accessible only to authenticated users."
    ]
    add_bullet_points(slide, points)

    # Slide 5: Fund Transfer Logic (Atomic)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Atomic Fund Transfers"
    points = [
        "SQL Transactions: Uses BEGIN/COMMIT to ensure data integrity.",
        "Balance Validation: Server-side check for sufficient funds.",
        "Recipient Verification: Real-time existence checks via database.",
        "Automated Logging: Instant transaction record generation."
    ]
    add_bullet_points(slide, points)

    # Slide 6: UI/UX Design System
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Premium User Experience"
    points = [
        "Glassmorphism: Translucent panels with backdrop blurring.",
        "Responsive Layout: Mobile-first approach with sidebars.",
        "Micro-animations: Pulsing badges, tray pop-ins, and focus effects.",
        "Dark Mode: Standard high-fidelity dark aesthetic for modern feel."
    ]
    add_bullet_points(slide, points)

    # Slide 7: Dashboard Hub
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Financial Dashboard"
    points = [
        "Account Summary: Detailed breakdown of balance and user info.",
        "Recent Activity: Scrollable transaction logs with status indicators.",
        "Product Management: Virtual Card system with mock details.",
        "Quick Actions: Easy access to transfers and settings."
    ]
    add_bullet_points(slide, points)

    # Slide 8: Deployment & Scalability
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Cloud Integration"
    points = [
        "GitHub-driven deployments for both frontend and backend.",
        "Database scalability with Neon Serverless PostgreSQL.",
        "CORS policy optimization for multi-cloud communication.",
        "Production-ready environment variables management."
    ]
    add_bullet_points(slide, points)

    # Slide 9: Future Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Future Enhancements"
    points = [
        "Implementing Real-time Multi-factor Authentication (MFA).",
        "Adding Real-time push notifications for transactions.",
        "Integrating PDF e-statement generation.",
        "Expanding to multi-currency and crypto support."
    ]
    add_bullet_points(slide, points)

    # Slide 10: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion"
    points = [
        "Nidhi Bank successfully blends security with modern aesthetics.",
        "Proven architecture using FastAPI and Next.js.",
        "Scalable database design ready for growth.",
        "A complete demonstration of full-stack engineering excellence."
    ]
    add_bullet_points(slide, points)

    # Save presentation
    prs.save("Nidhi_Bank_Project_Presentation.pptx")
    print("Presentation created successfully: Nidhi_Bank_Project_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
