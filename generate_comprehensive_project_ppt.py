import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def build_presentation():
    prs = Presentation()
    
    # helper functions for slide generation
    def set_dark_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42) # Midnight theme
        
    def add_custom_slide(title_text, content_list):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_dark_bg(slide)
        
        # Format title
        title_shape = slide.shapes.title
        title_shape.text = title_text
        if title_shape.has_text_frame:
            for p in title_shape.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(36)
                p.font.color.rgb = RGBColor(96, 165, 250)
                
        # Format content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        for point in content_list:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(241, 245, 249)
            p.space_after = Pt(12)
            
        return slide

    # Title Slide
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    set_dark_bg(s0)
    title = s0.shapes.title
    title.text = "NIDHI BANK: COMPLETE PROJECT DETAILS"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(96, 165, 250)
    
    subtitle = s0.placeholders[1]
    subtitle.text = "A Full-Stack Fintech Platform Architecture"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(226, 232, 240)

    # Slide Content Data
    slides_data = [
        ("1. Project Executive Summary", [
            "Project Name: Nidhi Bank",
            "Objective: To deliver a fully functional, secure, and modern digital banking system.",
            "Key Features: Real-time fund transfers, Virtual cards, User management, Security logging.",
            "Development Approach: Fully decoupled architecture with distinct Frontend, Backend, and DB."
        ]),
        ("2. Comprehensive Technology Stack", [
            "Frontend Ecosystem: Next.js 14, React.js, standard HTML/CSS structuring.",
            "Backend Framework: FastAPI (Python) for asynchronous, high-concurrency API handling.",
            "Database: Neon PostgreSQL, deployed as a serverless relational ledger.",
            "Hosting Providers: Vercel (Frontend Client) & Render (Backend API Service).",
            "Language Context: Python 3.x and JavaScript (ES6+)."
        ]),
        ("3. Project Directory Structure", [
            "/frontend: Client application encompassing Next.js components, pages, and CSS modules.",
            "/backend: Root for FastAPI application, python dependencies, and main.py route handler.",
            "/db: Database migration, seeding, and configuration scripts.",
            "Root files: Project deployment scripts (redeploy.py, check_render.py) and documentation."
        ]),
        ("4. Relational Database Modeling", [
            "Normalized Data Storage: Users and Transactions tables built with standard SQL schemas.",
            "Data Integrity: Foreign Key constraints linking Transactions directly to authentic user IDs.",
            "Numeric Precision: Using DECIMAL(12,2) for perfectly accurate non-floating financial balances.",
            "Auditing: Utilizing comprehensive TIMESTAMP logs to create immutable transaction records."
        ]),
        ("5. Secure Backend Architecture (FastAPI)", [
            "Authentication: Standard Bcrypt hashing avoiding older limit-bound libraries.",
            "Routing Logic: Segregated handlers for User operations over cross-origin authenticated requests.",
            "CORS Middleware: Custom CORS configuration enabling cross-platform interactions.",
            "Safety Measures: Real-time server-side balance checking prior to allowing any fund deduction."
        ]),
        ("6. Frontend Presentation & Navigation", [
            "Glassmorphism UI: Implementation of translucent UI elements with blurred aesthetic backdrops.",
            "Responsive Router: Next.js /app router utilized for handling deep linking (e.g. /dashboard/settings).",
            "Authentication States: Redirecting unauthenticated access and establishing localized auth persistence.",
            "Feedback Mechanisms: Professional toast notifications to replace standard browser alerts."
        ]),
        ("7. Foundational Features: Fund Transfers", [
            "End-to-End Flow: Frontend form validation linked directly to atomic backend SQL execution.",
            "Atomicity Guarantee: Wrapped in BEGIN/COMMIT blocks to prevent mismatched money movements.",
            "Validation: Strict system verification preventing users from transferring negative amounts or overdrafting.",
            "Instant Feedback: Web UI syncs with JSON responses immediately rendering new wallet balances."
        ]),
        ("8. Advanced Administrative Capabilities", [
            "Directory Listing: Database-connected UI listing hundreds of active users with fast filtering.",
            "Detailed Interactivity: Dropdown options providing dynamically generated virtual account details.",
            "Virtual Cards Ecosystem: Simulated infrastructure allowing users to generate functional card credentials."
        ]),
        ("9. Operational Deployment Pipeline", [
            "CI/CD Integration: Seamless triggers enabling auto-redeployments on Github master pushes.",
            "Vercel: Edge functional deployment providing extreme low latency for client files.",
            "Render: Scalable backend webservice operating on an automated deployment bridge.",
            "Automation Tools: Python scripts (e.g. check_vercel.py) automating daily dev-ops tests."
        ]),
        ("10. Conclusion & Project Finality", [
            "Outcome: A completely successful end-to-end banking implementation.",
            "Future Additions: Setup allows for effortless integration with mobile frameworks (e.g. React Native).",
            "Documentation: In-depth code comments and READMEs maintaining professional structure.",
            "Thank you for reviewing the NIDHI BANK comprehensive architecture."
        ])
    ]

    for title, points in slides_data:
        add_custom_slide(title, points)

    output_file = "NidhiBank_Comprehensive_Project_Details.pptx"
    prs.save(output_file)
    print(output_file)

if __name__ == "__main__":
    build_presentation()
