import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def start_ppt_build():
    prs = Presentation()

    # Paths to generated images (Absolute paths from brain dir)
    # Using forward slashes to avoid backslash issues
    brain_dir = "C:/Users/HP/.gemini/antigravity/brain/3fdaeae4-77f5-4960-a4fc-0bed68121c1f"
    hero_img = brain_dir + "/banking_hero_image_1775368895680.png"
    auth_img = brain_dir + "/secure_auth_icon_1775368915553.png"
    tech_img = brain_dir + "/tech_stack_visualization_1775368950930.png"

    def set_slide_background(slide, color=RGBColor(28, 28, 28)):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def apply_title_style(title_shape, color=RGBColor(0, 191, 255)):
        if not title_shape.has_text_frame:
            return
        text_frame = title_shape.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(36)
            paragraph.font.color.rgb = color

    def add_slide(layout_idx, title_text):
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide)
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            apply_title_style(slide.shapes.title)
        return slide

    def add_bullet_points(slide, points, font_size=Pt(20)):
        if len(slide.placeholders) < 2:
            return
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = font_size
            p.font.color.rgb = RGBColor(240, 240, 240)
            p.space_after = Pt(10)

    # --- Slide 0: Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    if os.path.exists(hero_img):
        slide.shapes.add_picture(hero_img, Inches(0), Inches(0), width=Inches(5))
    
    title_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(2))
    title_tf = title_box.text_frame
    title_p = title_tf.add_paragraph()
    title_p.text = "NIDHI BANK"
    title_p.font.bold = True
    title_p.font.size = Pt(54)
    title_p.font.color.rgb = RGBColor(0, 255, 255)
    
    sub_p = title_tf.add_paragraph()
    sub_p.text = "Next-Generation Digital Banking Infrastructure"
    sub_p.font.size = Pt(24)
    sub_p.font.color.rgb = RGBColor(200, 200, 200)

    # Slides 1-9
    content = [
        ("Executive Summary", [
            "Objective: To architect a high-security, ultra-responsive digital banking platform.",
            "Target Audience: Modern users seeking seamless fintech experiences.",
            "Core Value: Reliability, Security, and State-of-the-art Design.",
            "Implementation: Optimized Full-stack delivery via FastAPI & Next.js."
        ]),
        ("Technology Architecture", [
            "Frontend Interface: Next.js 14 utilizing CSR/SSR for optimal performance.",
            "Backend Engine: FastAPI with asynchronous PostgreSQL drivers.",
            "Database Engine: Serverless Neon PostgreSQL with high availability.",
            "Hosting Strategy: Vercel edge runtime & Render clusters."
        ]),
        ("System Functionality", [
            "Real-time Asset Management: Live balance tracking.",
            "Institutional Fund Transfers: Atomic SQL transactions.",
            "Personalized Banking: Custom virtual card generation.",
            "Administrative Governance: User auditing and risk analysis."
        ]),
        ("Security & Data Integrity", [
            "Cryptographic Hashing: Industry-standard Bcrypt encryption.",
            "CORS Enforcement: Deep-layer security against unauthorized requests.",
            "Input Sanitization: Robust validation to maintain data integrity.",
            "Session Persistence: Secure client-side state management."
        ]),
        ("Advanced UI/UX Philosophy", [
            "Design Trend: Modern Glassmorphism (Skeuomorphic evolution).",
            "Visual Cues: Backdrop blurs and depth perception.",
            "Accessibility: High-contrast typography and mobile-first logic.",
            "Engagement: Pulse animations and interactive tray models."
        ]),
        ("Relational Data Modeling", [
            "Normalization: Structured User and Transaction schemas.",
            "Atomicity: ACID compliant transfers ensuring consistency.",
            "Concurrency: Optimized connection pooling.",
            "Indexing: Strategic query optimization for history retrieval."
        ]),
        ("Technical Challenges Overcome", [
            "Issue: Asynchronous CORS race conditions. Solution: Middlewares.",
            "Issue: Dependency conflicts. Solution: Native Bcrypt implementation.",
            "Issue: Environment sync. Solution: Unified deployment scripts.",
            "Issue: Real-time UI sync. Solution: Optimized React hook patterns."
        ]),
        ("Innovation Roadmap", [
            "Q3 2026: Implementation of Multi-Factor Authentication (MFA).",
            "Q4 2026: Real-time fraud detection via heuristic analysis.",
            "2027: Integration with blockchain for decentralized auditing.",
            "Beyond: AI-driven financial advisory services."
        ]),
        ("Conclusion", [
            "NIDHI BANK delivers a professional, secure, and modern experience.",
            "Project demonstrates mastery of Full-Stack paradigms.",
            "The architecture is scalable and production-ready.",
            "Ready for real-world deployment and feature expansion."
        ])
    ]

    for title, points in content:
        s = add_slide(1, title)
        add_bullet_points(s, points)
        if title == "Technology Architecture" and os.path.exists(tech_img):
            s.shapes.add_picture(tech_img, Inches(6.5), Inches(4.5), width=Inches(3))
        if title == "Security & Data Integrity" and os.path.exists(auth_img):
            s.shapes.add_picture(auth_img, Inches(6.5), Inches(2), width=Inches(3))

    # --- Slide 10: Q&A ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide)
    if slide.shapes.title:
        slide.shapes.title.text = "THANK YOU"
        apply_title_style(slide.shapes.title, RGBColor(255, 215, 0))
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Questions & Discussions"
        slide.placeholders[1].font.color.rgb = RGBColor(255, 255, 255)

    prs.save("Nidhi_Bank_Professional_Presentation_Final.pptx")
    print("Professional presentation created successfully.")

if __name__ == "__main__":
    start_ppt_build()
