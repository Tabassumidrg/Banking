import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def build_detailed_ppt():
    prs = Presentation()
    brain_dir = "C:/Users/HP/.gemini/antigravity/brain/3fdaeae4-77f5-4960-a4fc-0bed68121c1f"
    hero_p = brain_dir + "/banking_hero_image_1775368895680.png"
    auth_p = brain_dir + "/secure_auth_icon_1775368915553.png"
    tech_p = brain_dir + "/tech_stack_visualization_1775368950930.png"

    def bg(s, c=RGBColor(15, 23, 42)): # Midnight Blue
        f = s.background.fill
        f.solid()
        f.fore_color.rgb = c

    def title_s(t, c=RGBColor(96, 165, 250)): # Sky Blue
        if t.has_text_frame:
            for p in t.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(36)
                p.font.color.rgb = c

    def add_s(idx, txt):
        s = prs.slides.add_slide(prs.slide_layouts[idx])
        bg(s)
        if s.shapes.title:
            s.shapes.title.text = txt
            title_s(s.shapes.title)
        return s

    def add_b(s, pts):
        if len(s.placeholders) < 2: return
        tf = s.placeholders[1].text_frame
        for pt in pts:
            p = tf.add_paragraph()
            p.text = pt
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(241, 245, 249)
            p.space_after = Pt(8)

    # 0. Title
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s0)
    if os.path.exists(hero_p): s0.shapes.add_picture(hero_p, Inches(0), Inches(0), width=Inches(5))
    tb = s0.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.5), Inches(3))
    t_p = tb.text_frame.add_paragraph()
    t_p.text = "NIDHI BANK: FULL DETAIL"
    t_p.font.bold = True
    t_p.font.size = Pt(48)
    t_p.font.color.rgb = RGBColor(59, 130, 246)
    s_p = tb.text_frame.add_paragraph()
    s_p.text = "Comprehensive Architectural Analysis"
    s_p.font.size = Pt(22)
    s_p.font.color.rgb = RGBColor(255,255,255)

    data = [
        ("1. System Architecture Diagram", [
            "Frontend: Next.js 14 (App Router) serving the client application.",
            "Backend: FastAPI providing a high-performance RESTful API layer.",
            "Database: Neon PostgreSQL (Serverless) for durable data storage.",
            "Hosting: Multi-provider deployment (Vercel for Frontend, Render for Backend).",
            "Communication: JSON-based REST APIs with CORS protection."
        ]),
        ("2. Frontend: Next.js Implementation", [
            "State Management: React 'useState' and 'useEffect' for local session handling.",
            "Client-Side Rendering: 'use client' directive for interactive dashboard elements.",
            "API Integration: Native 'fetch' API with async/await for non-blocking I/O.",
            "Session Persistence: Utilizing 'localStorage' for storing JWT-less user sessions.",
            "Real-time Updates: CustomEvent dispatchers for immediate UI feedback."
        ]),
        ("3. Backend: FastAPI Engineering", [
            "Framework: High-performance Python 3.x FastAPI (Asynchronous by design).",
            "Data Validation: Pydantic models (UserSignup, UserLogin) for schema enforcement.",
            "CORS Middleware: Explicit 'CORSMiddleware' configuration for global access.",
            "Dependency Management: 'psycopg2' with RealDictCursor for mapping SQL to JSON.",
            "Error Handling: Custom HTTPException wrappers for structured API responses."
        ]),
        ("4. Database: PostgreSQL Schema", [
            "Users Table: Primary identifier (SERIAL), Unique Email/Mobile, Hashed Passwords.",
            "Transactions Table: Foreign Keys linking to Users (sender_id, receiver_id).",
            "Data Types: 'DECIMAL(12,2)' for precision banking calculations.",
            "Timestamps: 'TIMESTAMP WITH TIME ZONE' for global transaction auditing.",
            "Normalization: 3NF design ensuring zero data redundancy."
        ]),
        ("5. Fund Transfer: Atomic Logic", [
            "Integrity: Wrapping balance updates in a single SQL TRANSACTION (BEGIN/COMMIT).",
            "Validation: Strict server-side verification of sender balance vs. transfer amount.",
            "Receiver Checks: Database-level existence verification before processing.",
            "Error Recovery: 'ROLLBACK' mechanism on any failure to prevent money loss.",
            "Logging: Automatic record insertion into the transactions table."
        ]),
        ("6. Security: Cryptographic Hardening", [
            "Hashing: Direct 'bcrypt' salt/hash implementation for credential storage.",
            "Prevention: Mitigation of 'passlib' vulnerabilities (72-character limit).",
            "CORS: Rigid Origin, Method, and Header white-listing for the Render API.",
            "Authentication: Secure sign-in flow with server-side credential verification.",
            "Input Sanitization: Mitigation of common SQL Injection and XSS vectors."
        ]),
        ("7. UI/UX: Design Philosophy", [
            "Theme: Midnight Blue Dark Mode using HSL CSS variables.",
            "Glassmorphism: High-fidelity 'backdrop-filter: blur(10px)' on all panels.",
            "Typography: 'Inter' sans-serif family for maximum readability.",
            "Animations: Keyframe-based fadeIn/Down/Up for a smooth product feel.",
            "Responsive: 100% Mobile-first UI using CSS Media Queries."
        ]),
        ("8. Administrative Management", [
            "User Hub: Advanced searchable list of all registered bank users.",
            "Dashboard: Real-time balance fetching with locale-aware (₹) formatting.",
            "Management tools: Dynamic transaction lists filtered by User ID.",
            "Virtual Cards: Logic-based card generation with mock security credentials.",
            "Audit Trail: Detailed logging of last 5 transactions on individual profiles."
        ]),
        ("9. SDLC & Deployment Pipeline", [
            "Version Control: Centralized GitHub repository architecture.",
            "CI/CD: Automated Git-triggers for Vercel (Frontend) and Render (Backend).",
            "Environment Sync: Centralized .env management across local/prod nodes.",
            "Build Verification: Rigid Next.js build checks to prevent runtime errors.",
            "API Redundancy: Multiple health-check endpoints for monitoring."
        ]),
        ("10. Future Innovation Roadmap", [
            "Authentication: Integration of Biometric and TOTP-based 2FA.",
            "Analytics: AI-driven spending patterns and risk score calculations.",
            "Internationalization: Support for multi-currency (USD, EUR) transfers.",
            "Architecture: Transitioning from monolith to modular Microservices.",
            "Scalability: Implementing Redis-based caching for high-traffic sessions."
        ]),
        ("11. Summary of Achievements", [
            "Fully functional end-to-end banking ecosystem built in < 1 week.",
            "Successfully mastered both Python (FastAPI) and JavaScript (Next.js) paradigms.",
            "High-end visual aesthetic on par with industry Fintech products.",
            "Scalable foundation ready for real-world user onboarding."
        ]),
        ("12. THANK YOU / Q&A", [
            "Technical Documentation: 100% coverage in PROJECT_TRACKER.md.",
            "Software Stack: Open source, modern, and high-performance.",
            "Presentation Author: Designed with Antigravity AI for perfection.",
            "Questions?", "Thank you for your time!"
        ])
    ]

    for t, p in data:
        st = add_s(1, t)
        add_b(st, p)
        if t == "1. System Architecture Diagram" and os.path.exists(tech_p): st.shapes.add_picture(tech_p, Inches(6.5), Inches(4.5), width=Inches(3))
        if t == "6. Security: Cryptographic Hardening" and os.path.exists(auth_p): st.shapes.add_picture(auth_p, Inches(6.8), Inches(2), width=Inches(2.5))

    prs.save("Nidhi_Bank_Full_Detail_Presentation.pptx")
    print("Full Detail PPT Done")

if __name__ == "__main__":
    build_detailed_ppt()
